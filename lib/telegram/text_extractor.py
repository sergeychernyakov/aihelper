import os
import docx
from bs4 import BeautifulSoup
from PyPDF2 import PdfReader
from pptx import Presentation
from textract import process
from striprtf.striprtf import rtf_to_text
import zipfile
import shutil
import subprocess
import pdfplumber

class TextExtractor:
    """
    A class to extract text from various file formats.
    Supported formats include txt, tex, docx, html, pdf, pptx, doc, rtfd, rtf, tar, zip.
    """
    ALLOWED_FILE_EXTENSIONS = {'.txt', '.tex', '.docx', '.html', '.pdf', '.pptx', '.doc', '.rtfd', '.rtf', '.tar', '.zip'}

    @staticmethod
    def extract_text(file_path):
        """
        Extracts text from a given file.
        :param file_path: str
            The path to the file from which text is to be extracted.
        :return: str
            Extracted text or a message indicating unsupported format.
        """
        ext = os.path.splitext(file_path)[1].lower()

        if ext in TextExtractor.ALLOWED_FILE_EXTENSIONS:
            if ext in ['.txt', '.tex']:
                return TextExtractor._read_plain_text(file_path)
            elif ext == '.docx':
                return TextExtractor._extract_from_docx(file_path)
            elif ext == '.html':
                return TextExtractor._extract_from_html(file_path)
            elif ext == '.pdf':
                return TextExtractor._extract_from_pdf(file_path)
            elif ext == '.pptx':
                return TextExtractor._extract_from_pptx(file_path)
            elif ext == '.doc':
                return TextExtractor._extract_from_doc(file_path)
            elif ext == '.rtf':
                return TextExtractor._extract_from_rtf(file_path)
            elif ext == '.pages':
                return TextExtractor._extract_from_pages(file_path)
            elif ext == '.rtfd':
                return TextExtractor._extract_from_rtfd(file_path)
            elif ext == '.tar':
                return "Tar file extraction not implemented"
            elif ext == '.zip':
                return "Zip file extraction not implemented"
        else:
            return "Unsupported file format"

    @staticmethod
    def _read_plain_text(file_path):
        """
        Reads plain text from a txt or tex file.
        :param file_path: str
            Path to the file.
        :return: str
            Content of the file.
        """
        with open(file_path, 'r') as file:
            return file.read()

    @staticmethod
    def _extract_from_docx(file_path):
        """
        Extracts text from a docx file.
        :param file_path: str
            Path to the docx file.
        :return: str
            Extracted text from the file.
        """
        doc = docx.Document(file_path)
        return '\n'.join(para.text for para in doc.paragraphs)

    @staticmethod
    def _extract_from_html(file_path):
        """
        Extracts text from an HTML file, ignoring script and style elements.
        :param file_path: str
            Path to the HTML file.
        :return: str
            Text content of the HTML file.
        """
        with open(file_path, 'r') as file:
            soup = BeautifulSoup(file, 'html.parser')

            # Remove script and style elements
            for script_or_style in soup(['script', 'style']):
                script_or_style.extract()

            # Get text
            text = soup.get_text()

            # Break into lines and remove leading and trailing space on each
            lines = (line.strip() for line in text.splitlines())

            # Break multi-headlines into a line each
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))

            # Drop blank lines
            text = '\n'.join(chunk for chunk in chunks if chunk)

            return text

    @staticmethod
    def _extract_from_pdf(file_path):
        """
        Extracts text from a PDF file using pdfplumber.
        :param file_path: str
            Path to the PDF file.
        :return: str
            Text content of the PDF file.
        """
        text = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text.append(page.extract_text())
        return '\n'.join(filter(None, text))

    @staticmethod
    def _extract_from_pptx(file_path):
        """
        Extracts text from a PPTX file.
        :param file_path: str
            Path to the PPTX file.
        :return: str
            Text content of the PPTX file.
        """
        ppt = Presentation(file_path)
        text = []
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)

    @staticmethod
    def _extract_from_doc(file_path):
        """
        Extracts text from a DOC file using textract.
        :param file_path: str
            Path to the DOC file.
        :return: str
            Text content of the DOC file.
        """
        try:
            return process(file_path).decode('utf-8')
        except Exception as e:
            return f"Error extracting text from .doc: {e}"

    @staticmethod
    def _extract_from_rtf(file_path):
        """
        Extracts text from an RTF file using the striprtf library.
        :param file_path: str
            Path to the RTF file.
        :return: str
            Text content of the RTF file.
        """
        try:
            with open(file_path, 'r') as file:
                rtf_content = file.read()
            return rtf_to_text(rtf_content)
        except Exception as e:
            return f"Error extracting text from .rtf: {e}"

    @staticmethod
    def _extract_from_rtfd(file_path):
        """
        Extracts text from an RTFD file.

        :param file_path: str
            Path to the RTFD file.
        :return: str
            Text content of the RTFD file.
        """
        temp_dir = os.path.join('tmp', 'rtfd_extraction')

        try:
            # Ensure the temporary directory exists
            os.makedirs(temp_dir, exist_ok=True)

            # Copy the RTFD (directory) content to the temporary directory
            shutil.copytree(file_path, temp_dir, dirs_exist_ok=True)

            # Search for an RTF file inside the directory
            for root, dirs, files in os.walk(temp_dir):
                for file in files:
                    if file.endswith('.rtf'):
                        rtf_file_path = os.path.join(root, file)
                        with open(rtf_file_path, 'r') as file:
                            # Extract and return the text content from the RTF file
                            return file.read()

            return "No RTF file found in RTFD package"
        except Exception as e:
            return f"Error extracting text from RTFD: {e}"
        finally:
            # Clean up the temporary directory
            shutil.rmtree(temp_dir, ignore_errors=True)


# Example usage:
# Suppose you have a document in one of the supported formats (e.g., 'document.docx').
# You can use the following code to extract text from it:

# {'.txt', '.tex', '.docx', '.html', '.pdf', '.pptx', '.doc', '.rtfd', '.rtf', '.tar', '.zip'}

# from lib.telegram.text_extractor import TextExtractor
# # Define a list of file paths with different extensions
# file_paths = [
# 'tmp/document.txt',
# 'tmp/document.tex',
# 'tmp/document.docx',
# 'tmp/document.html',
# 'tmp/document.pdf',
# 'tmp/document.pptx',
# 'tmp/document.doc',
# 'tmp/document.rtfd',
# 'tmp/document.rtf',
# # 'tmp/document.tar',  # Assuming tar and zip extraction not implemented yet
# # 'tmp/document.zip'
# ]


# file_paths = [
# 'tmp/document.pdf',
# ]

# # # Loop through each file path, extract text, and print it
# for file_path in file_paths:
#     print(f"\nExtracting text from {file_path}:")
#     extracted_text = TextExtractor.extract_text(file_path)
#     print(extracted_text)
